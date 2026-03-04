#!/usr/bin/env python3
"""Build a 5-minute, 5-slide deck using the Cisco Live template."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, List, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

VIDEO_URL = "https://www.youtube.com/watch?v=RnjgLlQTMf0"
TEMPLATE_PATH = Path(
    "/Users/rpias/Dropbox/new job 2025/presentation on ai/"
    "Cisco Live 2026 PPT Template_Dark_1772589711665001xsUU.pptx"
)
OUT_PATH = Path(
    "/Users/rpias/dev/vscode-dev-env/output/presentations/"
    "frontier-operations-cisco-template-5min.pptx"
)
SCRIPT_PATH = Path(
    "/Users/rpias/dev/vscode-dev-env/output/presentations/"
    "frontier-operations-cisco-template-5min-script.md"
)


def remove_all_slides(prs: Presentation) -> None:
    """Clear sample slides while preserving template masters/layouts."""
    slide_ids = list(prs.slides._sldIdLst)  # type: ignore[attr-defined]
    for slide_id in slide_ids:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)  # type: ignore[attr-defined]


def _find_title(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            return ph
    return None


def _find_bodies(slide) -> List:
    return sorted(
        [
            ph
            for ph in slide.placeholders
            if ph.placeholder_format.type == PP_PLACEHOLDER.BODY
        ],
        key=lambda p: p.left,
    )


def _find_table_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.TABLE:
            return ph
    return None


def set_title(slide, text: str) -> None:
    title = _find_title(slide)
    if title is None:
        raise RuntimeError("No title placeholder found on selected layout")
    title.text = text


def set_bullets(placeholder, lines: Sequence[Tuple[int, str]]) -> None:
    tf = placeholder.text_frame
    tf.clear()
    for i, (level, text) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level


def add_url_footer(prs: Presentation, slide) -> None:
    footer = slide.shapes.add_textbox(
        Inches(0.55), prs.slide_height - Inches(0.34), Inches(12.2), Inches(0.25)
    )
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Source video: {VIDEO_URL}"
    run.hyperlink.address = VIDEO_URL
    run.font.size = Pt(9)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(210, 210, 210)


def add_notes(slide, notes: str) -> None:
    slide.notes_slide.notes_text_frame.text = notes.strip()


def build_deck() -> None:
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE_PATH))
    remove_all_slides(prs)

    # Slide 1: Thesis
    s1 = prs.slides.add_slide(prs.slide_layouts[15])  # Title, 1 Column with Bullets
    set_title(s1, "Frontier Operations: The New AI Workforce Discipline")
    bodies = _find_bodies(s1)
    set_bullets(
        bodies[0],
        [
            (
                0,
                "Core thesis: business value now comes from operating at the moving boundary between human judgment and AI execution.",
            ),
            (0, "The speaker uses an expanding bubble model:"),
            (1, "Inside bubble = tasks agents can do reliably today."),
            (1, "Outside bubble = tasks still requiring human context, risk judgment, and accountability."),
            (
                0,
                "The boundary shifts every quarter, so skill durability depends on continuous recalibration, not one-time training.",
            ),
            (
                0,
                "This is why old workforce methods underperform in AI-native environments.",
            ),
        ],
    )
    add_url_footer(prs, s1)
    add_notes(
        s1,
        """
Slide one sets the frame. The video argues that we should stop treating AI capability as a fixed tool and start
treating it as a moving operating boundary. The bubble analogy is useful: inside the bubble are tasks agents can do
well right now, and outside are tasks where humans still create the most value. The high-leverage zone is the
surface, where delegation, verification, and intervention decisions happen. As models improve, tasks keep moving
inside the bubble, so yesterday's best process becomes today's bottleneck. The key implication is that career
resilience and company competitiveness now depend on recalibration speed. This is not a one-time literacy problem.
It is an ongoing operating discipline.
""",
    )

    # Slide 2: Five capabilities
    s2 = prs.slides.add_slide(prs.slide_layouts[16])  # Title, 2 Columns with Bullets
    set_title(s2, "The Five Capability Stack of Frontier Operations")
    bodies = _find_bodies(s2)
    set_bullets(
        bodies[0],
        [
            (0, "1) Boundary sensing"),
            (1, "Track where agents currently succeed and fail by task."),
            (0, "2) Seam design"),
            (1, "Engineer clean handoffs between agent and human phases."),
            (0, "3) Failure model maintenance"),
            (1, "Use task-specific checks, not generic skepticism."),
        ],
    )
    set_bullets(
        bodies[1],
        [
            (0, "4) Capability forecasting"),
            (1, "Make practical 6-12 month bets on what will automate next."),
            (0, "5) Leverage calibration"),
            (1, "Allocate scarce human attention by risk and business impact."),
            (0, "Operating principle"),
            (1, "Run all five continuously; this is a practice, not a checklist."),
        ],
    )
    add_url_footer(prs, s2)
    add_notes(
        s2,
        """
This slide summarizes the framework itself. First is boundary sensing: maintain a current map of what AI can and
cannot reliably do in your domain. Second is seam design: define exactly where work transitions between agents and
humans. Third is failure model maintenance: understand the specific ways outputs break for each task class. Fourth is
capability forecasting: make short-horizon bets so skills and workflows evolve ahead of the curve. Fifth is leverage
calibration: triage human attention to high-risk and high-value decisions instead of reviewing everything equally.
The speaker's important point is integration. Teams do not get leverage from isolated tactics. They get leverage by
running all five capabilities together, then updating them continuously as model behavior changes.
""",
    )

    # Slide 3: Organizational implications
    s3 = prs.slides.add_slide(prs.slide_layouts[15])  # Title, 1 Column with Bullets
    set_title(s3, "What Organizations Must Change Immediately")
    bodies = _find_bodies(s3)
    set_bullets(
        bodies[0],
        [
            (0, "Replace static AI training with practice environments and real delegation loops."),
            (
                0,
                "Measure calibration quality: can teams predict success/failure boundaries and choose correct checks?",
            ),
            (
                0,
                "Increase feedback density: many real cycles beat long one-time courses.",
            ),
            (
                0,
                "Create explicit frontier roles (automation leads / frontier operators).",
            ),
            (
                0,
                "Adopt leverage-oriented team structures: team-of-1 operators and team-of-5 pods with clear seam ownership.",
            ),
            (
                0,
                "Update hiring signals toward operational judgment, not generic prompting claims.",
            ),
        ],
    )
    add_url_footer(prs, s3)
    add_notes(
        s3,
        """
The operational takeaway is that AI adoption alone is not enough. Companies need a new system for how work evolves.
The speaker recommends practice-heavy environments, because capability intuition only develops through repeated real
delegation and verification cycles. Assessment should focus on calibration quality, not whether someone can write a
nice prompt. Organizationally, there should be explicit owners of frontier operations: people who keep failure models
current, redesign seams, and distribute process updates quickly. Team design also shifts. In some contexts, a single
high-skill operator can orchestrate substantial output. In broader product contexts, a small pod with one strong
frontier operator can perform like a much larger traditional team. This is a leverage model, not a headcount model.
""",
    )

    # Slide 4: 90-day plan
    s4 = prs.slides.add_slide(prs.slide_layouts[43])  # Title, Subtitle, Table
    set_title(s4, "90-Day Implementation Plan")
    bodies = _find_bodies(s4)
    if bodies:
        set_bullets(
            bodies[0],
            [
                (
                    0,
                    "Goal: establish a repeatable frontier operating system with measurable business impact.",
                )
            ],
        )
    table_ph = _find_table_placeholder(s4)
    if table_ph is None:
        raise RuntimeError("No table placeholder found on layout 43")
    table = table_ph.insert_table(rows=4, cols=3).table
    table.cell(0, 0).text = "Phase"
    table.cell(0, 1).text = "Execution Focus"
    table.cell(0, 2).text = "Deliverables / Metrics"

    table.cell(1, 0).text = "Days 1-30"
    table.cell(1, 1).text = "Map work + define seams + baseline current process"
    table.cell(1, 2).text = "Workflow map, risk tiers, baseline quality/cycle-time metrics"

    table.cell(2, 0).text = "Days 31-60"
    table.cell(2, 1).text = "Pilot team-of-1 and team-of-5 workflows"
    table.cell(2, 2).text = "Failure-check playbooks, attention rules, pilot throughput gains"

    table.cell(3, 0).text = "Days 61-90"
    table.cell(3, 1).text = "Scale + formal ownership + monthly recalibration"
    table.cell(3, 2).text = "Named frontier owners, KPI dashboard, rollout plan by function"

    add_url_footer(prs, s4)
    add_notes(
        s4,
        """
Here is a practical 90-day execution sequence. In the first month, map workflows, define where human-agent seams
sit, and baseline current quality and cycle time. In month two, run focused pilots using clear attention tiers and
task-specific failure checks. This is where teams learn what to automate, what to sample, and what always needs deep
review. In month three, formalize ownership and governance so the system survives beyond a pilot. Assign named
frontier operators, establish a monthly recalibration cadence, and track business outcomes on a simple dashboard.
Recommended KPIs are throughput per FTE, defect escape rate, decision cycle time, and the share of work that can run
agent-first at acceptable quality.
""",
    )

    # Slide 5: Summary + source
    s5 = prs.slides.add_slide(prs.slide_layouts[15])  # Title, 1 Column with Bullets
    set_title(s5, "Five-Minute Conclusion and Source Reference")
    bodies = _find_bodies(s5)
    set_bullets(
        bodies[0],
        [
            (
                0,
                "Bottom line: tools commoditize quickly; frontier operating skill does not.",
            ),
            (
                0,
                "Winning teams convert each model capability jump into reliable output faster than peers.",
            ),
            (
                0,
                "Start now: assign ownership, run pilots, and recalibrate every month.",
            ),
            (0, "Primary source video:"),
            (1, VIDEO_URL),
            (
                0,
                "Use this URL in follow-up materials to anchor team discussion and implementation planning.",
            ),
        ],
    )
    add_url_footer(prs, s5)
    add_notes(
        s5,
        """
To close, the strategic message from the video is straightforward. The differentiator is no longer access to models.
It is the human operational capability to convert model improvements into trustworthy business output. Teams that build
frontier operations early compound faster because they recalibrate continuously while others remain anchored to older
assumptions. The practical next step is not another generic AI workshop. It is explicit ownership, pilot workflows,
and a monthly rhythm of seam updates and failure-model refreshes. The source for this briefing is the linked YouTube
video shown on this slide and in the footer of each slide. If we apply the 90-day plan, we can move from AI adoption
to durable AI leverage.
""",
    )

    prs.save(str(OUT_PATH))


def write_script_file() -> None:
    script = f"""# Frontier Operations - 5 Minute Speaker Script

## Slide 1 - Frontier Operations: The New AI Workforce Discipline
The core message is that business value now comes from operating at a moving boundary between human judgment and AI execution. The speaker explains this with an expanding bubble model: inside the bubble are tasks agents can do reliably today, and outside are tasks that still need human context, accountability, and risk judgment. The highest-leverage work happens on the surface, where we decide what to delegate, how to verify, and when to intervene. Because model capability changes quickly, this boundary moves every quarter. That means old workforce approaches based on static training and certification are no longer enough. The new requirement is continuous recalibration.

## Slide 2 - The Five Capability Stack of Frontier Operations
The framework has five parts. Boundary sensing means maintaining a current map of where agents are strong and weak in your domain. Seam design means defining clean transitions between agent and human phases. Failure model maintenance means using task-specific checks instead of broad skepticism. Capability forecasting means making practical short-horizon bets about what will automate next. Leverage calibration means assigning human attention where risk and value are highest. The key point is integration: these five capabilities only deliver results when run together as an ongoing operating practice.

## Slide 3 - What Organizations Must Change Immediately
Organizations need to shift from static AI training to high-frequency practice environments. Skills improve through repeated real delegation and verification cycles, not one-time courses. Teams should be measured on calibration quality: can they predict where agents will succeed, where they may fail, and what checks are appropriate? Companies also need explicit frontier roles responsible for seam redesign, failure model updates, and process rollout. Team structure becomes leverage-first: some domains support high-output operators, while broader product work benefits from small pods with clear frontier ownership.

## Slide 4 - 90-Day Implementation Plan
In days 1 to 30, map workflows, define seams, and baseline quality and cycle-time metrics. In days 31 to 60, run focused pilots with clear attention tiers and targeted failure checks. In days 61 to 90, scale with formal ownership and monthly recalibration governance. The KPI set should stay simple and operational: throughput per FTE, defect escape rate, decision cycle time, and the percentage of work that can run agent-first at acceptable quality. The objective is to build a repeatable frontier operating system, not just a pilot demo.

## Slide 5 - Conclusion and Source
The conclusion is that tools commoditize, but frontier operating skill does not. The organizations that win are the ones that turn each model capability jump into reliable output faster than their peers. The immediate action is to assign ownership, run pilots, and establish a monthly recalibration cadence. Primary source video for this briefing: {VIDEO_URL}. Use that source link in follow-up discussions so teams stay anchored to the framework and can move directly into execution planning.
"""
    SCRIPT_PATH.write_text(script, encoding="utf-8")


def main() -> None:
    build_deck()
    write_script_file()
    print(OUT_PATH)
    print(SCRIPT_PATH)


if __name__ == "__main__":
    main()
