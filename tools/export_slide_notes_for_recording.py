#!/usr/bin/env python3
"""Export speaker notes from each slide into per-slide text files for recording."""

from __future__ import annotations

import argparse
from pathlib import Path
import re

from pptx import Presentation


def sanitize(text: str) -> str:
    cleaned = text.replace("\r", "\n")
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
    return cleaned


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Export slide speaker notes to text files for voice recording."
    )
    parser.add_argument("--pptx", required=True, help="Path to source presentation")
    parser.add_argument(
        "--out-dir",
        required=True,
        help="Output directory for slideXX.txt note files",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    pptx_path = Path(args.pptx).expanduser().resolve()
    out_dir = Path(args.out_dir).expanduser().resolve()

    if not pptx_path.exists():
        raise SystemExit(f"Error: PPTX not found: {pptx_path}")

    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(pptx_path))

    index_lines = []
    for i, slide in enumerate(prs.slides, start=1):
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text or ""
        notes = sanitize(notes)
        if not notes:
            notes = "[No speaker notes found for this slide.]"

        txt_path = out_dir / f"slide{i:02d}.txt"
        txt_path.write_text(notes + "\n", encoding="utf-8")
        word_count = len([w for w in notes.split() if w.strip()])
        index_lines.append(f"slide{i:02d}.txt\twords={word_count}")

    (out_dir / "index.txt").write_text("\n".join(index_lines) + "\n", encoding="utf-8")
    print(f"Wrote {len(prs.slides)} note files to {out_dir}")


if __name__ == "__main__":
    main()
