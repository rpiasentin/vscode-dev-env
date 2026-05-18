#!/usr/bin/env python3
"""Create an editable PPTX deck from AI datacenter sales org-chart artifacts."""

from __future__ import annotations

import argparse
import csv
import json
import math
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
DEFAULT_BUNDLE_ROOT = REPO_ROOT / "output" / "research" / "cisco-org-ai-datacenter-sales"
DEFAULT_OUTPUT = DEFAULT_BUNDLE_ROOT / "ai-datacenter-sales-org-charts.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(17, 34, 64)
INK = RGBColor(32, 42, 55)
MUTED = RGBColor(96, 111, 132)
PAPER = RGBColor(249, 247, 241)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(166, 180, 199)
GOLD = RGBColor(245, 190, 72)
GOLD_PALE = RGBColor(255, 245, 214)
BLUE = RGBColor(59, 105, 180)
BLUE_PALE = RGBColor(229, 239, 255)
GREEN = RGBColor(39, 135, 101)
GREEN_PALE = RGBColor(224, 246, 237)
RED = RGBColor(171, 64, 64)


def load_json(path: Path) -> Dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def read_team_rows(path: Path) -> List[Dict[str, str]]:
    with path.open(newline="", encoding="utf-8") as fh:
        return list(csv.DictReader(fh))


def sorted_team_rows(rows: Sequence[Dict[str, str]]) -> List[Dict[str, str]]:
    return sorted(
        rows,
        key=lambda row: (
            -int(row.get("candidate_count") or 0),
            row.get("db_label", ""),
            row.get("manager_alias", ""),
        ),
    )


def set_slide_background(slide: Any, color: RGBColor = PAPER) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_run_style(run: Any, *, size: int, bold: bool = False, color: RGBColor = INK, font: str = "Aptos") -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = INK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, size=size, bold=bold, color=color)
    return box


def add_multiline_text(
    slide: Any,
    lines: Sequence[Tuple[str, int, bool, RGBColor]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for index, (text, size, bold, color) in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(1)
        run = paragraph.add_run()
        run.text = text
        set_run_style(run, size=size, bold=bold, color=color)
    return box


def add_card(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = CARD,
    line: RGBColor = LINE,
    radius_shape: Any = MSO_SHAPE.ROUNDED_RECTANGLE,
    line_width: float = 1.0,
) -> Any:
    shape = slide.shapes.add_shape(radius_shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def add_metric_card(slide: Any, label: str, value: str, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    add_card(slide, x, y, w, h, fill=CARD, line=color, line_width=1.3)
    add_text(slide, value, x + 0.12, y + 0.12, w - 0.24, 0.32, size=22, bold=True, color=color)
    add_text(slide, label, x + 0.12, y + 0.51, w - 0.24, 0.28, size=8, bold=True, color=MUTED)


def add_footer(slide: Any, slide_no: int) -> None:
    add_text(
        slide,
        "Local Cisco directory data | generated on this Mac",
        0.42,
        7.18,
        5.5,
        0.18,
        size=7,
        color=MUTED,
    )
    add_text(slide, str(slide_no), 12.45, 7.18, 0.42, 0.18, size=7, color=MUTED, align=PP_ALIGN.RIGHT)


def node_label(node: Dict[str, Any], *, compact: bool = False) -> List[Tuple[str, int, bool, RGBColor]]:
    name = node.get("full_name") or node.get("alias") or "Unknown"
    alias = node.get("alias") or ""
    title = node.get("title") or ""
    if compact:
        if len(title) > 58:
            title = title[:55] + "..."
        return [
            (name, 7, True, INK),
            (title, 5, False, MUTED),
            (alias, 5, False, MUTED),
        ]
    if len(title) > 80:
        title = title[:77] + "..."
    return [
        (name, 9, True, INK),
        (title, 6, False, MUTED),
        (alias, 6, False, MUTED),
    ]


def find_node(nodes_by_alias: Dict[str, Dict[str, Any]], alias: str) -> Dict[str, Any]:
    return nodes_by_alias.get(alias, {"alias": alias, "full_name": alias, "title": ""})


def edge_children(chart: Dict[str, Any], parent_alias: str) -> List[str]:
    children = [
        edge["child_alias"]
        for edge in chart.get("edges", [])
        if edge.get("parent_alias") == parent_alias
    ]
    return sorted(set(children))


def ancestor_path(chart: Dict[str, Any], manager_alias: str) -> List[str]:
    parent_by_child = {
        edge["child_alias"]: edge["parent_alias"]
        for edge in chart.get("edges", [])
        if edge.get("child_alias") != edge.get("parent_alias")
    }
    path: List[str] = []
    current = manager_alias
    seen: set[str] = set()
    while current in parent_by_child and current not in seen:
        seen.add(current)
        parent = parent_by_child[current]
        path.append(parent)
        current = parent
    return list(reversed(path))


def add_line(slide: Any, x1: float, y1: float, x2: float, y2: float, *, color: RGBColor = LINE, width: float = 1.0) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)


def add_node_card(
    slide: Any,
    node: Dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    selected: bool = False,
    manager: bool = False,
    compact: bool = False,
) -> None:
    fill = GOLD_PALE if selected else BLUE_PALE if manager else CARD
    line = GOLD if selected else BLUE if manager else LINE
    add_card(slide, x, y, w, h, fill=fill, line=line, line_width=1.4 if selected or manager else 0.8)
    add_multiline_text(slide, node_label(node, compact=compact), x + 0.05, y + 0.03, w - 0.1, h - 0.06)


def add_summary_slide(prs: Presentation, manifest: Dict[str, Any], team_rows: Sequence[Dict[str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_text(slide, "AI Datacenter Infrastructure Sales Org Charts", 0.45, 0.35, 9.8, 0.45, size=23, bold=True, color=NAVY)
    add_text(slide, "Focused local org charts built from completed Oliver and Jeetu crawls", 0.48, 0.82, 8.7, 0.28, size=10, color=MUTED)

    add_metric_card(slide, "MATCHED PEOPLE", str(manifest.get("candidate_count", 0)), 0.55, 1.45, 2.15, 0.95, BLUE)
    add_metric_card(slide, "TEAM CHARTS", str(manifest.get("team_count", 0)), 2.95, 1.45, 2.15, 0.95, GREEN)
    add_metric_card(slide, "LARGEST CHART", str(max((int(row["node_count"]) for row in team_rows), default=0)), 5.35, 1.45, 2.15, 0.95, GOLD)
    add_metric_card(slide, "WARNINGS", str(sum(int(row["warning_count"]) for row in team_rows)), 7.75, 1.45, 2.15, 0.95, RED)

    add_card(slide, 0.55, 2.85, 5.8, 3.4, fill=CARD, line=LINE)
    add_text(slide, "Selection Rule", 0.8, 3.08, 2.4, 0.26, size=13, bold=True, color=NAVY)
    add_multiline_text(
        slide,
        [
            ("Each included person matched all three marker families:", 10, False, INK),
            ("1. Sales / GTM", 10, True, BLUE),
            ("2. AI / accelerated compute", 10, True, GOLD),
            ("3. Datacenter / infrastructure", 10, True, GREEN),
            ("Charts are grouped by immediate manager, with peer/direct-report context and the leadership path up to the focused root.", 8, False, MUTED),
        ],
        0.78,
        3.48,
        5.35,
        2.45,
    )

    add_card(slide, 6.7, 2.85, 5.7, 3.4, fill=CARD, line=LINE)
    add_text(slide, "Database Coverage", 6.95, 3.08, 2.6, 0.26, size=13, bold=True, color=NAVY)
    y = 3.55
    summaries = manifest.get("summaries", {})
    for db_label in sorted(summaries):
        summary = summaries[db_label]
        add_metric_card(
            slide,
            f"{db_label.upper()} MATCHES",
            str(summary.get("candidate_count", 0)),
            7.0,
            y,
            2.15,
            0.82,
            BLUE if db_label == "oliver" else GREEN,
        )
        add_metric_card(
            slide,
            f"{db_label.upper()} TEAMS",
            str(summary.get("team_count", 0)),
            9.45,
            y,
            2.15,
            0.82,
            GREEN if db_label == "oliver" else BLUE,
        )
        y += 1.05
    add_footer(slide, 1)


def add_index_slide(prs: Presentation, team_rows: Sequence[Dict[str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_text(slide, "Team Chart Index", 0.45, 0.35, 5.5, 0.45, size=23, bold=True, color=NAVY)
    add_text(slide, "Sorted by number of matched people in the team", 0.48, 0.82, 6.2, 0.28, size=10, color=MUTED)

    rows = sorted_team_rows(team_rows)
    cols = 2
    per_col = math.ceil(len(rows) / cols)
    start_slide = 3
    for col in range(cols):
        x = 0.55 + col * 6.25
        for index, row in enumerate(rows[col * per_col : (col + 1) * per_col]):
            absolute_index = col * per_col + index
            y = 1.35 + index * 0.58
            add_card(slide, x, y, 5.75, 0.42, fill=CARD, line=LINE, line_width=0.5)
            add_text(slide, f"{start_slide + absolute_index:02d}", x + 0.12, y + 0.11, 0.35, 0.12, size=7, bold=True, color=MUTED)
            add_text(slide, row["manager_alias"], x + 0.55, y + 0.07, 2.3, 0.15, size=8, bold=True, color=INK)
            add_text(
                slide,
                f"{row['candidate_count']} matches | {row['direct_report_count']} direct reports | {row['node_count']} nodes",
                x + 0.55,
                y + 0.24,
                4.65,
                0.12,
                size=6,
                color=MUTED,
            )
    add_footer(slide, 2)


def add_team_slide(prs: Presentation, chart_payload: Dict[str, Any], team_row: Dict[str, str], slide_no: int) -> None:
    chart = chart_payload["chart"]
    manager_alias = chart_payload["manager_alias"]
    nodes_by_alias = {node["alias"]: node for node in chart.get("nodes", [])}
    children = edge_children(chart, manager_alias)
    selected_aliases = set(chart.get("selected_aliases", []))
    manager_node = find_node(nodes_by_alias, manager_alias)
    ancestors = [alias for alias in ancestor_path(chart, manager_alias) if alias != manager_alias]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    title_name = manager_node.get("full_name") or manager_alias
    add_text(slide, title_name, 0.45, 0.26, 8.9, 0.38, size=20, bold=True, color=NAVY)
    add_text(
        slide,
        f"{chart['db_label']} | manager alias {manager_alias} | {team_row['candidate_count']} matched people",
        0.48,
        0.68,
        8.4,
        0.24,
        size=8,
        color=MUTED,
    )
    add_metric_card(slide, "MATCHES", team_row["candidate_count"], 9.35, 0.25, 1.2, 0.65, GOLD)
    add_metric_card(slide, "DIRECTS", team_row["direct_report_count"], 10.7, 0.25, 1.2, 0.65, BLUE)
    add_metric_card(slide, "NODES", team_row["node_count"], 12.05, 0.25, 0.85, 0.65, GREEN)

    add_text(slide, "Leadership path", 0.55, 1.13, 2.2, 0.2, size=9, bold=True, color=NAVY)
    rail_x, rail_y, rail_w = 0.52, 1.42, 2.42
    rail_h = 4.95
    add_card(slide, rail_x, rail_y, rail_w, rail_h, fill=RGBColor(253, 252, 248), line=LINE, line_width=0.7)
    if ancestors:
        card_h = min(0.63, max(0.42, (rail_h - 0.35) / max(len(ancestors), 1) - 0.08))
        total_h = len(ancestors) * card_h + max(len(ancestors) - 1, 0) * 0.08
        y = rail_y + max(0.18, (rail_h - total_h) / 2)
        previous_mid: Optional[Tuple[float, float]] = None
        for alias in ancestors:
            add_node_card(slide, find_node(nodes_by_alias, alias), rail_x + 0.18, y, rail_w - 0.36, card_h, compact=True)
            mid = (rail_x + rail_w / 2, y + card_h)
            if previous_mid:
                add_line(slide, previous_mid[0], previous_mid[1], rail_x + rail_w / 2, y, width=0.8)
            previous_mid = mid
            y += card_h + 0.08
    else:
        add_text(slide, "No upstream chain available", rail_x + 0.25, rail_y + 0.45, rail_w - 0.5, 0.24, size=8, color=MUTED)

    add_text(slide, "Team chart", 3.25, 1.13, 2.2, 0.2, size=9, bold=True, color=NAVY)
    manager_x, manager_y, manager_w, manager_h = 5.65, 1.18, 3.05, 0.76
    add_node_card(slide, manager_node, manager_x, manager_y, manager_w, manager_h, manager=True)
    if ancestors:
        add_line(slide, rail_x + rail_w, rail_y + rail_h / 2, manager_x, manager_y + manager_h / 2, width=1.1)

    grid_x, grid_y, grid_w, grid_h = 3.25, 2.34, 9.55, 4.3
    add_card(slide, grid_x - 0.05, grid_y - 0.14, grid_w + 0.1, grid_h + 0.22, fill=RGBColor(253, 252, 248), line=LINE, line_width=0.6)
    add_text(slide, "Direct reports and matched sellers", grid_x + 0.1, grid_y - 0.04, 4.4, 0.18, size=8, bold=True, color=MUTED)

    if children:
        cols = 4 if len(children) > 12 else 3 if len(children) > 6 else max(1, min(len(children), 3))
        rows = math.ceil(len(children) / cols)
        gutter_x = 0.12
        gutter_y = 0.13
        card_w = (grid_w - (cols - 1) * gutter_x) / cols
        card_h = min(0.68, (grid_h - 0.36 - (rows - 1) * gutter_y) / rows)
        start_y = grid_y + 0.32
        trunk_y = manager_y + manager_h + 0.18
        add_line(slide, manager_x + manager_w / 2, manager_y + manager_h, manager_x + manager_w / 2, trunk_y, width=0.8)
        for index, alias in enumerate(children):
            row = index // cols
            col = index % cols
            x = grid_x + col * (card_w + gutter_x)
            y = start_y + row * (card_h + gutter_y)
            selected = alias in selected_aliases
            add_node_card(slide, find_node(nodes_by_alias, alias), x, y, card_w, card_h, selected=selected, compact=True)
            add_line(slide, manager_x + manager_w / 2, trunk_y, x + card_w / 2, y, color=GOLD if selected else LINE, width=0.65)
    else:
        add_text(slide, "No direct reports available in chart payload", grid_x + 0.25, grid_y + 0.55, 4.0, 0.3, size=9, color=MUTED)

    legend_y = 6.82
    add_card(slide, 3.25, legend_y, 0.18, 0.18, fill=GOLD_PALE, line=GOLD)
    add_text(slide, "Matched AI datacenter sales signal", 3.5, legend_y - 0.01, 2.7, 0.18, size=7, color=MUTED)
    add_card(slide, 6.25, legend_y, 0.18, 0.18, fill=BLUE_PALE, line=BLUE)
    add_text(slide, "Immediate manager", 6.5, legend_y - 0.01, 1.7, 0.18, size=7, color=MUTED)
    if chart.get("warnings"):
        add_text(slide, f"{len(chart['warnings'])} chart warning(s)", 8.45, legend_y - 0.01, 2.0, 0.18, size=7, color=RED)
    add_footer(slide, slide_no)


def build_presentation(bundle_root: Path, output_path: Path) -> Path:
    manifest = load_json(bundle_root / "manifest.json")
    team_rows = sorted_team_rows(read_team_rows(bundle_root / "teams.csv"))
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_summary_slide(prs, manifest, team_rows)
    add_index_slide(prs, team_rows)
    for slide_index, row in enumerate(team_rows, start=3):
        chart_payload = load_json(Path(row["json_path"]))
        add_team_slide(prs, chart_payload, row, slide_index)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--bundle-root", type=Path, default=DEFAULT_BUNDLE_ROOT)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    output = build_presentation(args.bundle_root, args.output)
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
