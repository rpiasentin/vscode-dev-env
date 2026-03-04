#!/usr/bin/env python3
"""Build a 5-slide information-dense deck from the Frontier Operations video summary."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(24, 49, 83)
HEADER_BG = RGBColor(232, 240, 252)
TEXT_COLOR = RGBColor(35, 35, 35)


def set_run_style(run, size=20, bold=False, color=TEXT_COLOR):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_title(slide, title, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_style(r, size=34, bold=True, color=TITLE_COLOR)
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        stb = slide.shapes.add_textbox(
            Inches(0.55), Inches(0.95), Inches(12.0), Inches(0.5)
        )
        stf = stb.text_frame
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        set_run_style(sr, size=14, bold=False, color=RGBColor(70, 70, 70))


def add_bullets(slide, x, y, w, h, lines, font_size=18, level_map=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0 if level_map is None else level_map[i]
        p.space_after = Pt(8)
        p.font.size = Pt(font_size)
        p.font.name = "Calibri"
        p.font.color.rgb = TEXT_COLOR
    return box


def build_deck(out_file: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Frontier Operations: The New Core Workforce Skill",
        "Source: YouTube video RnjgLlQTMf0 | Theme: human-AI boundary as a moving operating layer",
    )

    left_shape = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.4), Inches(6.25), Inches(5.7)
    )
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = RGBColor(247, 250, 255)
    left_shape.line.color.rgb = RGBColor(201, 215, 235)
    add_bullets(
        slide,
        0.75,
        1.7,
        5.8,
        5.2,
        [
            "Core claim: unlike legacy skills, AI-era value comes from operating at a shifting boundary.",
            "Bubble model:",
            "Inside bubble = tasks AI agents can perform reliably today.",
            "Outside bubble = work still requiring human context, risk judgment, and accountability.",
            "Surface = high-value frontier where delegation, verification, and intervention decisions happen.",
            "As models improve, tasks migrate inward; workers must continuously recalibrate.",
        ],
        font_size=17,
        level_map=[0, 0, 1, 1, 1, 0],
    )

    right_shape = slide.shapes.add_shape(
        1, Inches(6.95), Inches(1.4), Inches(5.9), Inches(5.7)
    )
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = RGBColor(252, 247, 239)
    right_shape.line.color.rgb = RGBColor(228, 207, 171)
    add_bullets(
        slide,
        7.2,
        1.7,
        5.35,
        5.1,
        [
            "Why this matters now:",
            "Boundary skill has no fixed finish line; refresh cycle is near-quarterly.",
            "Operational gap: training systems still optimize for static certification models.",
            "Resulting risk: over-trust, underuse, and outdated verification habits.",
            "Opportunity: frontier operators create leverage by redesigning work faster than peers.",
        ],
        font_size=17,
        level_map=[0, 1, 1, 1, 1],
    )

    slide.notes_slide.notes_text_frame.text = (
        "Today the main claim is that workforce skills used to have a finish line, "
        "but AI changes that pattern. Imagine AI capability as an expanding bubble. "
        "Inside are tasks agents can do reliably, outside are tasks still needing humans. "
        "The highest-value zone is the boundary between the two, where delegation, "
        "verification, and intervention decisions happen. As models improve, tasks move inward "
        "and the boundary shifts. Career resilience now depends on continuously staying calibrated "
        "to that moving edge."
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "The Five Frontier Operations Capabilities",
        "Integrated, continuous, and mutually reinforcing",
    )

    table = slide.shapes.add_table(6, 4, Inches(0.45), Inches(1.4), Inches(12.4), Inches(5.7)).table
    table.columns[0].width = Inches(2.25)
    table.columns[1].width = Inches(3.15)
    table.columns[2].width = Inches(3.35)
    table.columns[3].width = Inches(3.65)

    headers = ["Capability", "What It Means", "Common Failure If Missing", "Operator Action"]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG

    rows = [
        (
            "1) Boundary sensing",
            "Continuously update where AI succeeds/fails in your domain.",
            "Using last quarter assumptions; over-trust or under-delegation.",
            "Run recurring calibration tasks after major model/tool updates.",
        ),
        (
            "2) Seam design",
            "Engineer clean handoffs between agent and human phases.",
            "Messy transitions, hidden rework, unclear ownership.",
            "Define artifacts and checks at each transition boundary.",
        ),
        (
            "3) Failure model maintenance",
            "Track task-specific failure patterns, not generic skepticism.",
            "Either review everything or miss subtle high-risk errors.",
            "Map failure mode by task type and attach targeted verification.",
        ),
        (
            "4) Capability forecasting",
            "Make 6-12 month bets on what moves inside AI territory next.",
            "Chasing every tool or waiting too long to adapt workflows.",
            "Shift learning investment toward rising leverage layers.",
        ),
        (
            "5) Leverage calibration",
            "Allocate scarce human attention by risk and value tier.",
            "Uniform review depth leads to bottlenecks or blind trust.",
            "Use triage thresholds: automate routine, review risk-critical paths.",
        ),
    ]

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, text in enumerate(row):
            table.cell(r_idx, c_idx).text = text

    for r in range(6):
        for c in range(4):
            tf = table.cell(r, c).text_frame
            for p in tf.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(13 if r else 14)
                p.font.bold = bool(r == 0)
                p.font.color.rgb = TEXT_COLOR

    slide.notes_slide.notes_text_frame.text = (
        "The framework has five capabilities: boundary sensing, seam design, failure model maintenance, "
        "capability forecasting, and leverage calibration. These are not checklist items. "
        "They run simultaneously. High performance comes from integrating all five at once and "
        "updating them continuously as model behavior changes."
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Operating Model for Organizations",
        "From static training to continuous calibration systems",
    )

    add_bullets(
        slide,
        0.65,
        1.55,
        6.2,
        2.8,
        [
            "Design principles:",
            "Build practice environments (real tasks, realistic failures, changing rules).",
            "Measure calibration quality, not just prompt fluency.",
            "Maximize feedback density from daily delegation plus review loops.",
            "Create explicit frontier roles (workflow redesign plus failure governance).",
        ],
        font_size=17,
        level_map=[0, 1, 1, 1, 1],
    )

    ops_table = slide.shapes.add_table(3, 3, Inches(0.65), Inches(4.05), Inches(12.05), Inches(2.75)).table
    ops_table.columns[0].width = Inches(2.0)
    ops_table.columns[1].width = Inches(5.1)
    ops_table.columns[2].width = Inches(4.95)

    for c, h in enumerate(["Structure", "When It Works Best", "Leverage Pattern"]):
        ops_table.cell(0, c).text = h
        ops_table.cell(0, c).fill.solid()
        ops_table.cell(0, c).fill.fore_color.rgb = HEADER_BG

    ops_table.cell(1, 0).text = "Team of 1"
    ops_table.cell(1, 1).text = "Strong operator plus narrow domain plus tight feedback loops"
    ops_table.cell(1, 2).text = "1 person orchestrates multi-agent workflows; 5-10x legacy output"

    ops_table.cell(2, 0).text = "Team of 5"
    ops_table.cell(2, 1).text = "1 deep frontier operator plus specialists plus execution cadence"
    ops_table.cell(2, 2).text = "Pod ships like a much larger team via seam and attention design"

    for r in range(3):
        for c in range(3):
            tf = ops_table.cell(r, c).text_frame
            for p in tf.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(13 if r else 14)
                p.font.bold = bool(r == 0)
                p.font.color.rgb = TEXT_COLOR

    slide.notes_slide.notes_text_frame.text = (
        "Organizations should replace static AI training with dynamic operations: practice environments, "
        "calibration-based assessment, high feedback density, and explicit frontier roles. "
        "Output increasingly scales with leverage quality rather than headcount alone. "
        "Team-of-1 and team-of-5 structures are practical templates."
    )

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "Hiring and Management Playbook",
        "Assess operational judgment, not credential proxies",
    )

    card1 = slide.shapes.add_shape(1, Inches(0.55), Inches(1.45), Inches(6.15), Inches(5.85))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(241, 250, 243)
    card1.line.color.rgb = RGBColor(166, 209, 174)
    add_bullets(
        slide,
        0.85,
        1.75,
        5.7,
        5.3,
        [
            "Strong signals to hire or promote:",
            "Can map current AI success and failure boundaries in specific domain work.",
            "Can redesign workflow seams when a new capability appears.",
            "Maintains differentiated failure checks by task class.",
            "Shows repeatable short-horizon forecasting accuracy.",
            "Can explain attention triage policy across risk tiers.",
        ],
        font_size=16,
        level_map=[0, 1, 1, 1, 1, 1],
    )

    card2 = slide.shapes.add_shape(1, Inches(6.95), Inches(1.45), Inches(5.85), Inches(5.85))
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(255, 247, 247)
    card2.line.color.rgb = RGBColor(222, 173, 173)
    add_bullets(
        slide,
        7.25,
        1.75,
        5.35,
        5.3,
        [
            "Anti-patterns to avoid:",
            "Generic 'good at prompting' claims without operational evidence.",
            "Uniform review depth across all agent outputs.",
            "No named owner for evolving human-agent boundary design.",
            "Management controls:",
            "Define what is auto-approved, sampled, and always human-reviewed.",
            "Review calibration drift monthly and update verification protocols.",
        ],
        font_size=16,
        level_map=[0, 1, 1, 1, 0, 1, 1],
    )

    slide.notes_slide.notes_text_frame.text = (
        "Better hiring and promotion signals are calibration-in-context and redesign capability. "
        "Managers need explicit attention policies describing what is automated, sampled, and deeply reviewed. "
        "This avoids both review bottlenecks and blind trust."
    )

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(
        slide,
        "90-Day Implementation Roadmap",
        "Translate Frontier Operations into measurable business outcomes",
    )

    roadmap = slide.shapes.add_table(4, 3, Inches(0.55), Inches(1.45), Inches(8.2), Inches(4.4)).table
    roadmap.columns[0].width = Inches(1.7)
    roadmap.columns[1].width = Inches(2.6)
    roadmap.columns[2].width = Inches(3.9)
    for c, h in enumerate(["Phase", "Focus", "Deliverables"]):
        roadmap.cell(0, c).text = h
        roadmap.cell(0, c).fill.solid()
        roadmap.cell(0, c).fill.fore_color.rgb = HEADER_BG

    roadmap.cell(1, 0).text = "Days 1-30"
    roadmap.cell(1, 1).text = "Map and Baseline"
    roadmap.cell(1, 2).text = "Workflow map, seam definitions, quality and cycle-time baseline"

    roadmap.cell(2, 0).text = "Days 31-60"
    roadmap.cell(2, 1).text = "Pilot Pods"
    roadmap.cell(2, 2).text = "Team pilots, attention tiers, targeted failure-check playbooks"

    roadmap.cell(3, 0).text = "Days 61-90"
    roadmap.cell(3, 1).text = "Scale and Govern"
    roadmap.cell(3, 2).text = "Named frontier owners, monthly recalibration cadence, wider rollout"

    for r in range(4):
        for c in range(3):
            tf = roadmap.cell(r, c).text_frame
            for p in tf.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(13 if r else 14)
                p.font.bold = bool(r == 0)
                p.font.color.rgb = TEXT_COLOR

    add_bullets(
        slide,
        9.0,
        1.6,
        4.0,
        2.6,
        [
            "Core KPIs:",
            "Throughput per FTE",
            "Defect escape rate",
            "Decision cycle time",
            "Agent-first coverage with acceptable quality",
        ],
        font_size=15,
        level_map=[0, 1, 1, 1, 1],
    )

    add_bullets(
        slide,
        9.0,
        4.15,
        4.0,
        2.8,
        [
            "Executive takeaway:",
            "Tools commoditize; frontier operating skill does not.",
            "Winning teams turn model upgrades into reliable output faster.",
            "Assign ownership now; unmanaged boundary drift is a strategic risk.",
        ],
        font_size=15,
        level_map=[0, 1, 1, 1],
    )

    slide.notes_slide.notes_text_frame.text = (
        "Use a 90-day sequence: map and baseline, pilot with explicit review tiers, then scale with ownership and governance. "
        "Track throughput, defects, cycle time, and agent-first quality. "
        "Strategic point: tools alone are not the moat; operational capability is."
    )

    prs.save(out_file)


def main():
    out_dir = Path("/Users/rpias/dev/vscode-dev-env/output/presentations")
    out_dir.mkdir(parents=True, exist_ok=True)
    out_file = out_dir / "frontier-operations-5-slide-deck.pptx"
    build_deck(out_file)
    print(out_file)


if __name__ == "__main__":
    main()
